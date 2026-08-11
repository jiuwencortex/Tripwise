"""
Generate BLOGPOST-KeyNote-Huawei.pptx
Huawei Developer Conference style:
  - Primary red: #c7000b
  - Dark background: #0a0a10
  - Secondary: #e87722 (OpenJiuwen orange)
  - Professional, structured, enterprise tone
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Theme colours
# ---------------------------------------------------------------------------
DARK    = RGBColor(0x0a, 0x0a, 0x10)
DARK2   = RGBColor(0x16, 0x16, 0x24)
RED     = RGBColor(0xc7, 0x00, 0x0b)
ORANGE  = RGBColor(0xe8, 0x77, 0x22)
WHITE   = RGBColor(0xf0, 0xf0, 0xf0)
LGRAY   = RGBColor(0xa0, 0xa0, 0xb0)
DRED    = RGBColor(0x7a, 0x00, 0x07)
GREEN   = RGBColor(0x00, 0xb0, 0x50)
BLUE    = RGBColor(0x00, 0x70, 0xc0)
TEAL    = RGBColor(0x2A, 0x9D, 0x8F)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_slide(bg=DARK):
    slide = prs.slides.add_slide(BLANK)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def txt(slide, text, x, y, w, h, size=20, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Calibri"
    return tb


def add_multiline(slide, lines, x, y, w, line_height=Inches(0.45),
                  size=18, color=WHITE, bold=False, bullet="•  "):
    for i, line in enumerate(lines):
        txt(slide, f"{bullet}{line}", x, y + i * line_height, w,
            line_height + Inches(0.05),
            size=size, color=color, bold=bold, align=PP_ALIGN.LEFT)


def rect(slide, x, y, w, h, fill, text="", tsize=15, tcolor=WHITE,
         bold=False, border_color=None, border_w=1):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border_color:
        sh.line.color.rgb = border_color
        sh.line.width = Pt(border_w)
    else:
        sh.line.fill.background()
    if text:
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(tsize)
        r.font.color.rgb = tcolor
        r.font.bold = bold
        r.font.name = "Calibri"
        from pptx.util import Pt as PPt
        tf.margin_top    = PPt(4)
        tf.margin_bottom = PPt(4)
        tf.margin_left   = PPt(8)
        tf.margin_right  = PPt(8)
    return sh


def red_bar(slide, y=Inches(0.0), height=Inches(0.12)):
    """Huawei-style top red accent bar."""
    rect(slide, 0, y, W, height, RED)


def section_header(slide, title, sub=""):
    """Standard section header layout."""
    red_bar(slide)
    txt(slide, title, Inches(0.7), Inches(0.3), Inches(11.9), Inches(0.75),
        size=30, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    if sub:
        txt(slide, sub, Inches(0.7), Inches(1.05), Inches(11.9), Inches(0.45),
            size=16, color=LGRAY, italic=True, align=PP_ALIGN.LEFT)
    # thin divider
    d = slide.shapes.add_shape(1, Inches(0.7), Inches(1.55), Inches(11.9), Pt(1))
    d.fill.solid(); d.fill.fore_color.rgb = RED; d.line.fill.background()


def arrow_right(slide, x, y, length, color=LGRAY):
    c = slide.shapes.add_connector(1, x, y, x + length, y)
    c.line.color.rgb = color
    c.line.width = Pt(2)


def arrow_down(slide, x, y, length, color=LGRAY):
    c = slide.shapes.add_connector(1, x, y, x, y + length)
    c.line.color.rgb = color
    c.line.width = Pt(2)


def center_x(w): return (W - w) / 2


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = add_slide()
# Full-width red top bar (thicker for title)
rect(s, 0, 0, W, Inches(0.2), RED)
# Huawei-style diagonal red accent bottom strip
rect(s, 0, Inches(6.9), W, Inches(0.6), RED)

txt(s, "OpenJiuwen Community  ·  JiuwenClaw Platform", Inches(0.7), Inches(0.4),
    Inches(11.9), Inches(0.5), size=14, color=LGRAY, align=PP_ALIGN.LEFT)

txt(s, "Building the Intelligent Foundation",
    Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.1),
    size=46, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
txt(s, "for the Next Decade",
    Inches(0.7), Inches(2.55), Inches(11.9), Inches(1.0),
    size=46, color=RED, bold=True, align=PP_ALIGN.LEFT)

txt(s, "JiuwenClaw: Architect AND Intelligence Engine",
    Inches(0.7), Inches(3.75), Inches(11.9), Inches(0.55),
    size=20, color=LGRAY, italic=True, align=PP_ALIGN.LEFT)

txt(s, "Huawei Developer Conference  ·  2026",
    Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.4),
    size=14, color=WHITE, align=PP_ALIGN.LEFT)

# ---------------------------------------------------------------------------
# SLIDE 2 — One Principle
# ---------------------------------------------------------------------------
s = add_slide()
section_header(s, "One Architectural Principle",
               "That will define intelligent systems for the decade ahead")

# Large quote box
rect(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(2.1),
     DARK2, border_color=RED, border_w=2)
txt(s, "The agent that creates the thin client\nis also the intelligence backbone the thin client calls at runtime.",
    Inches(1.0), Inches(1.9), Inches(11.3), Inches(1.3),
    size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Three boxes below
roles = [
    ("🏗️  Architect",        "Designs and engineers\nthe thin client",         RED),
    ("⚡  Intelligence Engine","Powers all reasoning when\nthe product calls it", ORANGE),
    ("🔄  Self-Evolves",      "Skills Autonomous Evolution —\ncontinuous improvement", TEAL),
]
bw = Inches(3.7); bh = Inches(1.5); by = Inches(4.0)
gap = Inches(0.3)
total = 3 * bw + 2 * gap
bx_start = (W - total) / 2
for i, (label, desc, col) in enumerate(roles):
    bx = bx_start + i * (bw + gap)
    rect(s, bx, by, bw, bh, col,
         f"{label}\n{desc}", tsize=15, tcolor=WHITE, bold=True)

txt(s, "One agent. Two roles. Continuous evolution.",
    Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.5),
    size=18, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# SLIDE 3 — The Two-Role Architecture
# ---------------------------------------------------------------------------
s = add_slide()
section_header(s, "The Two-Role Architecture",
               "JiuwenClaw and TripWise — a working demonstration")

# Left: text explanation
txt(s, "Phase 1: Creation", Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.5),
    size=18, color=RED, bold=True)
add_multiline(s, [
    "JiuwenClaw engineers the thin client",
    "Backend · Frontend · Prompts · Schemas",
    "Complete architecture — zero domain logic in client",
], Inches(0.7), Inches(2.25), Inches(5.5), size=15, color=WHITE,
   line_height=Inches(0.5))

txt(s, "Phase 2: Runtime", Inches(0.7), Inches(3.85), Inches(5.5), Inches(0.5),
    size=18, color=ORANGE, bold=True)
add_multiline(s, [
    "Thin client calls JiuwenClaw on every action",
    "JiuwenClaw: multi-step reasoning, tools, memory",
    "Returns schema-validated JSON — rendered by client",
], Inches(0.7), Inches(4.4), Inches(5.5), size=15, color=WHITE,
   line_height=Inches(0.5))

# Right: diagram
jcx = Inches(7.5); jcy = Inches(2.2); jcw = Inches(4.5); jch = Inches(1.0)
rect(s, jcx, jcy, jcw, jch, RED,
     "🧠 JiuwenClaw\nArchitect AND Intelligence Engine",
     tsize=15, tcolor=WHITE, bold=True)

twx = Inches(7.5); twy = Inches(4.8); tww = Inches(4.5); twh = Inches(1.0)
rect(s, twx, twy, tww, twh, DARK2,
     "🖥️ TripWise\n(Thin Client — No Business Logic)",
     tsize=15, tcolor=WHITE, bold=False,
     border_color=ORANGE, border_w=2)

mid_x = jcx + jcw / 2
# Arrow: JC → TW
a1 = s.shapes.add_connector(1, mid_x, jcy + jch,
                              mid_x, twy)
a1.line.color.rgb = RED; a1.line.width = Pt(2)
txt(s, "① Built", mid_x + Inches(0.1), Inches(3.65), Inches(1.5), Inches(0.4),
    size=13, color=RED, bold=True)

# Arrow: TW → JC (offset slightly to the right)
mid_x2 = mid_x + Inches(0.5)
a2 = s.shapes.add_connector(1, mid_x2, twy,
                              mid_x2, jcy + jch)
a2.line.color.rgb = ORANGE; a2.line.width = Pt(2)
txt(s, "② Calls for intelligence", mid_x2 + Inches(0.1), Inches(3.75),
    Inches(2.5), Inches(0.35), size=12, color=ORANGE, italic=True)

# ---------------------------------------------------------------------------
# SLIDE 4 — JiuwenClaw Capabilities
# ---------------------------------------------------------------------------
s = add_slide()
section_header(s, "JiuwenClaw: Full-Stack Intelligent Agent Infrastructure")

caps = [
    ("Architect",            "Designs and engineers complete thin-client applications\n— from architecture to prompt schemas to UI",                                        RED),
    ("Power",                "Provides multi-step reasoning, tool use, and browser\nautomation when called by the products it built",                                      ORANGE),
    ("Evolve",               "Skills Autonomous Evolution — captures failures\nand self-improves continuously without redeployment",                                       TEAL),
    ("Integrate",            "WebSocket (JiuwenClaw) and REST (OpenJiuwen) interfaces\nfor seamless integration into any platform",                                        BLUE),
]
bw = Inches(5.5); bh = Inches(1.15)
by_start = Inches(1.75); gap = Inches(0.2)
for i, (cap, desc, col) in enumerate(caps):
    by2 = by_start + i * (bh + gap)
    rect(s, Inches(0.7), by2, Inches(1.5), bh, col, cap,
         tsize=17, tcolor=WHITE, bold=True)
    rect(s, Inches(2.3), by2, Inches(10.3), bh, DARK2, desc,
         tsize=15, tcolor=WHITE, bold=False,
         border_color=RGBColor(0x2a, 0x2a, 0x3e), border_w=1)

txt(s, "The same agent that engineers the thin client is the intelligence the thin client calls.",
    Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5),
    size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# SLIDE 5 — 2026 Trend Convergence
# ---------------------------------------------------------------------------
s = add_slide()
section_header(s, "2026 Trend Convergence",
               "Eight industry trends — all pointing at one architecture")

trends = [
    ("Harnessing Agents",        "Thin clients call JiuwenClaw via structured prompts + context"),
    ("Agent Substrate",          "Remove JiuwenClaw and the application is an empty shell"),
    ("Self-Refining Systems",    "Skills Autonomous Evolution — improves without human redeployment"),
    ("Ambient Agents",           "Operates real browsers: cookies, sessions, authenticated interfaces"),
    ("Executable Intelligence",  "Schema-validated JSON, directly parsed into application state"),
    ("Business-Native Agents",   "JiuwenClaw IS the domain logic — built by agent, called by thin client"),
    ("Agentic Operating Layers", "Applications call JiuwenClaw — the intelligence layer between app and system"),
    ("Autonomous Verticalization","One architecture. Any domain. Zero structural changes."),
]
by_start = Inches(1.75)
row_h = Inches(0.57)
for i, (trend, impl) in enumerate(trends):
    by2 = by_start + i * (row_h + Inches(0.03))
    col = RED if i % 2 == 0 else DRED
    rect(s, Inches(0.7), by2, Inches(3.3), row_h, col,
         trend, tsize=12, tcolor=WHITE, bold=True)
    rect(s, Inches(4.1), by2, Inches(8.9), row_h, DARK2,
         impl, tsize=12, tcolor=WHITE,
         border_color=RGBColor(0x2a, 0x2a, 0x3e))

# ---------------------------------------------------------------------------
# SLIDE 6 — TripWise: Scenario Demonstration
# ---------------------------------------------------------------------------
s = add_slide()
section_header(s, "Scenario Demonstration: TripWise",
               "Zero domain logic in the client — all intelligence called from JiuwenClaw")

txt(s, "Key insight:",
    Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.4),
    size=15, color=RED, bold=True)
rect(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(0.75),
     DARK2, "TripWise is a pure thin client. It has no domain logic. "
            "JiuwenClaw built it. TripWise calls JiuwenClaw at runtime for every answer.",
     tsize=15, tcolor=WHITE, border_color=RED, border_w=1)

# 6-step flow
steps = ["👤 User\nInput", "🌍 Destinations", "✈️ Flights", "🏨 Hotels", "🎭 Attractions", "📅 Itinerary\n+ Map"]
bw2 = Inches(1.65); bh2 = Inches(0.9); by3 = Inches(3.1)
gap2 = Inches(0.22)
total2 = 6 * bw2 + 5 * gap2
sx = (W - total2) / 2
cols = [DARK2, RED, RED, RED, RED, RED]
for i, (step, col) in enumerate(zip(steps, cols)):
    bx2 = sx + i * (bw2 + gap2)
    r = rect(s, bx2, by3, bw2, bh2, col, step, tsize=13, tcolor=WHITE, bold=(col == RED))
    if col == RED:
        r.line.color.rgb = ORANGE; r.line.width = Pt(1)
    if i < len(steps) - 1:
        ax = bx2 + bw2; ay = by3 + bh2 / 2
        a = s.shapes.add_connector(1, ax, ay, ax + gap2, ay)
        a.line.color.rgb = LGRAY; a.line.width = Pt(1.5)

# Call annotation below flow
txt(s, "Each → = TripWise calls JiuwenClaw  ·  JiuwenClaw returns JSON  ·  TripWise renders",
    Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.45),
    size=13, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

# Ecosystem header
txt(s, "The openJiuwen Ecosystem — Two Equal Integration Paths:",
    Inches(0.7), Inches(4.55), Inches(11.9), Inches(0.4),
    size=14, color=WHITE, bold=True)

# JiuwenClaw box
rect(s, Inches(0.7), Inches(5.0), Inches(5.7), Inches(0.8),
     RED, "🧠 JiuwenClaw\nLive hosted agent  ·  WebSocket  ·  No setup needed",
     tsize=13, tcolor=WHITE, bold=True)
txt(s, "Use the live community agent", Inches(0.7), Inches(5.85), Inches(5.7), Inches(0.35),
    size=11, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

# OpenJiuwen box
rect(s, Inches(6.7), Inches(5.0), Inches(5.9), Inches(0.8),
     DARK2, "🔧 OpenJiuwen\nBuild your own private agent  ·  REST  ·  Full ownership",
     tsize=13, tcolor=WHITE, border_color=ORANGE, border_w=2)
txt(s, "Own your intelligence layer", Inches(6.7), Inches(5.85), Inches(5.9), Inches(0.35),
    size=11, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

# Ecosystem principle
txt(s, "Both paths are architecturally identical from TripWise's perspective. Same prompts. Same JSON back. The ecosystem answers your way.",
    Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.45),
    size=12, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# SLIDE 7 — Universal Pattern
# ---------------------------------------------------------------------------
s = add_slide()
section_header(s, "The Pattern Is Universal",
               "Any domain where intelligence is needed")

domains = [
    ("Travel & Hospitality",   "Trip planning · Booking · Concierge"),
    ("Enterprise Operations",  "Procurement · HR · Compliance"),
    ("Healthcare",             "Patient intake · Diagnostics · Care pathways"),
    ("Financial Services",     "Loan processing · Portfolio · Risk"),
    ("Smart Devices",          "HarmonyOS · Cross-device · Ambient computing"),
]
dw = Inches(5.4); dh = Inches(0.8)
for i, (domain, items) in enumerate(domains):
    col = i % 2
    row = i // 2 if i < 4 else 2
    by_d = Inches(1.75) + row * (dh + Inches(0.3))
    bx_d = Inches(0.7) + col * (dw + Inches(0.5))
    if i == 4:
        bx_d = center_x(dw)
    rect(s, bx_d, by_d, dw, dh, DARK2,
         f"{domain}\n{items}", tsize=14, tcolor=WHITE,
         border_color=RED, border_w=1)

txt(s, "One agent architecture. Any domain. Built by JiuwenClaw. Powered by JiuwenClaw.",
    Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5),
    size=17, color=RED, bold=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# SLIDE 8 — Developer + Enterprise Opportunity
# ---------------------------------------------------------------------------
s = add_slide()
section_header(s, "The Opportunity", "For developers and enterprises")

# Left column — developers
txt(s, "For Developers", Inches(0.7), Inches(1.75), Inches(5.8), Inches(0.5),
    size=20, color=RED, bold=True)
dev_points = [
    "No hardcoded logic — JiuwenClaw holds it all",
    "Change behavior by changing the prompt",
    "Unified intelligence across all devices + HarmonyOS",
    "Open-source — full transparency, full control",
]
add_multiline(s, dev_points, Inches(0.7), Inches(2.35), Inches(5.8),
              size=16, color=WHITE, line_height=Inches(0.58))

# Right column — enterprise
txt(s, "For Enterprises", Inches(7.2), Inches(1.75), Inches(5.8), Inches(0.5),
    size=20, color=ORANGE, bold=True)
ent_points = [
    "Entire business domains built in weeks",
    "Workflows adapt — change the prompt, thin client follows",
    "Self-improvement with no redeployment",
    "One agent: architect AND intelligence backbone",
]
add_multiline(s, ent_points, Inches(7.2), Inches(2.35), Inches(5.8),
              size=16, color=WHITE, line_height=Inches(0.58), bullet="→  ")

# Bottom statement
rect(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.8),
     RED,
     "By 2030, intelligent orchestration will be foundational for every digital enterprise.\n"
     "JiuwenClaw positions organizations to lead that transition today.",
     tsize=15, tcolor=WHITE, bold=False)

# ---------------------------------------------------------------------------
# SLIDE 9 — Closing
# ---------------------------------------------------------------------------
s = add_slide()
red_bar(s, y=Inches(0), height=Inches(0.15))
rect(s, 0, Inches(6.85), W, Inches(0.65), RED)

txt(s, "The architect and the intelligence engine",
    Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.9),
    size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "are the same agent.",
    Inches(0.7), Inches(2.1), Inches(11.9), Inches(0.9),
    size=52, color=RED, bold=True, align=PP_ALIGN.CENTER)

# Three-box summary
bw3 = Inches(3.5); bh3 = Inches(1.0); by_b = Inches(3.3)
gap3 = Inches(0.27)
total3 = 3 * bw3 + 2 * gap3
bx3_start = (W - total3) / 2
for i, (label, col) in enumerate([
    ("🏗️  JiuwenClaw\nBuilds the Thin Client",  DRED),
    ("⚡  Thin Client\nCalls JiuwenClaw",         DARK2),
    ("🔄  JiuwenClaw\nSelf-Evolves",              TEAL),
]):
    bx3 = bx3_start + i * (bw3 + gap3)
    r3 = rect(s, bx3, by_b, bw3, bh3, col, label, tsize=15, tcolor=WHITE, bold=True)
    if col == DARK2:
        r3.line.color.rgb = ORANGE; r3.line.width = Pt(2)
    if i < 2:
        ax3 = bx3 + bw3; ay3 = by_b + bh3 / 2
        a3 = s.shapes.add_connector(1, ax3, ay3, ax3 + gap3, ay3)
        a3.line.color.rgb = WHITE; a3.line.width = Pt(2)

txt(s, "This is the architecture of the next decade.",
    Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.5),
    size=20, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Together, we will build it.",
    Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.5),
    size=20, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)

txt(s, "🌐 openjiuwen.com/en   ·   💻 github.com/openJiuwen-ai   ·   🧠 gitcode.com/openJiuwen",
    Inches(0.7), Inches(7.0), Inches(11.9), Inches(0.4),
    size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = "/Users/mishka/PycharmProjects/openjiuwen/michael-agent-studio/backend/openjiuwen_studio/core/dsl_converter/demos/TripwisePython/blog/jiuwenclaw_first/BLOGPOST-KeyNote-Huawei.pptx"
prs.save(out)
print(f"Saved: {out}")
