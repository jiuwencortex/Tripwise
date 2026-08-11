"""
Generate BLOGPOST-KeyNote-SteveJobs.pptx from the Steve Jobs keynote script.
Apple dark aesthetic: #1d1d1f bg, #f5f5f7 text, #0071e3 accent.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy

# ---------------------------------------------------------------------------
# Theme colours
# ---------------------------------------------------------------------------
DARK    = RGBColor(0x1d, 0x1d, 0x1f)
WHITE   = RGBColor(0xf5, 0xf5, 0xf7)
BLUE    = RGBColor(0x00, 0x71, 0xe3)
ORANGE  = RGBColor(0xE7, 0x6F, 0x51)
GREEN   = RGBColor(0x34, 0xc7, 0x59)
GRAY    = RGBColor(0x86, 0x86, 0x8b)
TEAL    = RGBColor(0x2A, 0x9D, 0x8F)
RED     = RGBColor(0xE7, 0x2C, 0x2C)

W = Inches(13.33)   # 16:9 slide width
H = Inches(7.5)     # 16:9 slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_text(slide, text, x, y, w, h, size=24, color=WHITE, bold=False,
             align=PP_ALIGN.CENTER, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Helvetica Neue"
    return txBox


def add_rect(slide, x, y, w, h, fill_color, text="", txt_size=16,
             txt_color=WHITE, bold=False, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(txt_size)
        run.font.color.rgb = txt_color
        run.font.bold = bold
        run.font.name = "Helvetica Neue"
        from pptx.util import Pt as PPt
        tf.margin_top    = PPt(6)
        tf.margin_bottom = PPt(6)
        tf.margin_left   = PPt(8)
        tf.margin_right  = PPt(8)
    return shape


def add_arrow_right(slide, x, y, length, color=GRAY):
    """Draw a simple horizontal right-pointing connector."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        x, y, x + length, y
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def add_arrow_down(slide, x, y, length, color=GRAY):
    connector = slide.shapes.add_connector(
        1,
        x, y, x, y + length
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def center_x(width):
    return (W - width) / 2


def slide_subtitle(slide, text, y=Inches(1.15), color=GRAY, size=18):
    add_text(slide, text, Inches(0.5), y, Inches(12.33), Inches(0.6),
             size=size, color=color, align=PP_ALIGN.CENTER, italic=True)


def divider_line(slide, y=Inches(1.0), color=RGBColor(0x3a, 0x3a, 0x3c)):
    shape = slide.shapes.add_shape(1, Inches(1), y, Inches(11.33), Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "There's Something in the Air…",
         Inches(1), Inches(1.8), Inches(11.33), Inches(1.2),
         size=44, color=WHITE, bold=False, align=PP_ALIGN.CENTER)
add_text(s, "Intelligence.",
         Inches(1), Inches(3.0), Inches(11.33), Inches(1.2),
         size=64, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "JiuwenClaw  ·  openJiuwen Community  ·  2026",
         Inches(1), Inches(6.5), Inches(11.33), Inches(0.6),
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 2 — The Old World
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "The Old World", Inches(1), Inches(0.5), Inches(11.33), Inches(0.8),
         size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.4))
lines = [
    "Engineers write logic.",
    "Systems execute logic.",
    "Two separate things. Two separate teams.",
    "Six months to ship.",
    "The world changes.",
    "You do it all over again.",
]
for i, line in enumerate(lines):
    add_text(s, line, Inches(2.5), Inches(1.6 + i * 0.72), Inches(8.33), Inches(0.65),
             size=22, color=WHITE if i < 4 else GRAY, align=PP_ALIGN.CENTER)
add_text(s, "There's got to be a better way.",
         Inches(1), Inches(6.5), Inches(11.33), Inches(0.65),
         size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 3 — A New Pattern
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "A New Pattern", Inches(1), Inches(0.5), Inches(11.33), Inches(0.8),
         size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.4))

rows = [
    ("From apps with logic",          "→",  "to apps that delegate logic"),
    ("From systems that execute",      "→",  "to systems that generate instructions"),
    ("From software that stays fixed", "→",  "to software that evolves itself"),
]
for i, (left, arrow, right) in enumerate(rows):
    y = Inches(1.8 + i * 1.3)
    add_text(s, left,  Inches(0.5), y, Inches(4.8), Inches(0.9), size=20, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(s, arrow, Inches(5.4), y, Inches(0.8), Inches(0.9), size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, right, Inches(6.3), y, Inches(6.5), Inches(0.9), size=20, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

add_text(s, "This is a new architecture.",
         Inches(1), Inches(6.4), Inches(11.33), Inches(0.7),
         size=26, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 4 — Introducing JiuwenClaw
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "Introducing JiuwenClaw",
         Inches(1), Inches(0.3), Inches(11.33), Inches(0.8),
         size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.2))

# Central JiuwenClaw box
cx = Inches(4.9); cy = Inches(1.5); cw = Inches(3.5); ch = Inches(0.95)
add_rect(s, cx, cy, cw, ch, DARK,
         "🧠 JiuwenClaw\n— Architect AND Intelligence Engine —",
         txt_size=15, txt_color=WHITE, bold=True)
# Add a visible border to the JiuwenClaw box
jc_shape = s.shapes[-1]
jc_shape.line.color.rgb = BLUE
jc_shape.line.width = Pt(2)

# Left branch — Builds
add_rect(s, Inches(1.2), Inches(3.2), Inches(3.5), Inches(0.85),
         ORANGE, "🏗️ Builds the Product\nCode · Prompts · Schemas", txt_size=14)
add_arrow_down(s, Inches(2.95), Inches(2.45), Inches(0.75), ORANGE)

# Right branch — Powers (called at runtime)
add_rect(s, Inches(8.4), Inches(3.2), Inches(3.5), Inches(0.85),
         BLUE, "⚡ Powers Intelligence\n(Called at runtime)", txt_size=14)
add_arrow_down(s, Inches(10.15), Inches(2.45), Inches(0.75), BLUE)

# Connecting line from JC to left
connector = s.shapes.add_connector(1, Inches(4.9), Inches(1.97),
                                    Inches(2.95), Inches(3.2))
connector.line.color.rgb = ORANGE
connector.line.width = Pt(2)
# Connecting line from JC to right
connector2 = s.shapes.add_connector(1, Inches(8.4), Inches(1.97),
                                     Inches(10.15), Inches(3.2))
connector2.line.color.rgb = BLUE
connector2.line.width = Pt(2)

# Product box at bottom — calls JiuwenClaw for intelligence
add_rect(s, Inches(4.4), Inches(5.1), Inches(4.5), Inches(0.85),
         RGBColor(0x00, 0x40, 0x80), "🖥️ The Product\n(Calls JiuwenClaw for all intelligence)",
         txt_size=13, txt_color=WHITE)
# arrow from Build branch to product (JC built it)
s.shapes.add_connector(1, Inches(2.95), Inches(4.05), Inches(5.3), Inches(5.1)).line.color.rgb = ORANGE
# arrow from product back to Intelligence branch (product calls JC)
s.shapes.add_connector(1, Inches(7.5), Inches(5.1), Inches(10.15), Inches(4.05)).line.color.rgb = BLUE

add_text(s, "The same agent that built the product is the intelligence the product calls.",
         Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.7),
         size=18, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 5 — The Reveal: TripWise
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "Let Me Show You",
         Inches(1), Inches(0.4), Inches(11.33), Inches(0.8),
         size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.25))

add_text(s, "TripWise", Inches(1), Inches(1.5), Inches(11.33), Inches(0.9),
         size=52, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "A travel planner.", Inches(1), Inches(2.5), Inches(11.33), Inches(0.6),
         size=22, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "TripWise contains", Inches(1), Inches(3.3), Inches(11.33), Inches(0.6),
         size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "ZERO", Inches(1), Inches(3.9), Inches(11.33), Inches(0.9),
         size=72, color=RED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "travel logic.", Inches(1), Inches(4.85), Inches(11.33), Inches(0.6),
         size=24, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "JiuwenClaw built it.   TripWise calls it for every decision.   Same agent.",
         Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.7),
         size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 6 — Watch What Happens (demo flow)
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "Watch What Happens",
         Inches(1), Inches(0.4), Inches(11.33), Inches(0.7),
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.15))

steps = [
    ("👤", "You\n30 sec"),
    ("🌍", "Destinations"),
    ("✈️", "Flights"),
    ("🏨", "Hotels"),
    ("🎭", "Attractions"),
    ("📅", "Itinerary\n+ Map"),
]
box_w = Inches(1.55)
box_h = Inches(1.35)
gap   = Inches(0.28)
total = len(steps) * box_w + (len(steps) - 1) * gap
start_x = (W - total) / 2
y_box = Inches(2.2)

colors = [DARK, RGBColor(0x1a, 0x4a, 0x2e), RGBColor(0x0a, 0x3a, 0x6e),
          RGBColor(0x1a, 0x3a, 0x5e), RGBColor(0x3a, 0x1a, 0x5e),
          BLUE]
border_colors = [GRAY, GREEN, BLUE, RGBColor(0x2A, 0x9D, 0x8F),
                 RGBColor(0x9B, 0x5D, 0xE5), WHITE]

for i, (icon, label) in enumerate(steps):
    bx = start_x + i * (box_w + gap)
    add_rect(s, bx, y_box, box_w, box_h, colors[i],
             f"{icon}\n{label}", txt_size=15, txt_color=WHITE, bold=(i == len(steps)-1))
    sh = s.shapes[-1]
    sh.line.color.rgb = border_colors[i]
    sh.line.width = Pt(1.5)
    if i < len(steps) - 1:
        ax = bx + box_w
        ay = y_box + box_h / 2
        arr = s.shapes.add_connector(1, ax, ay, ax + gap, ay)
        arr.line.color.rgb = GRAY
        arr.line.width = Pt(2)

add_text(s, "Each step: JiuwenClaw reasons from everything decided so far.",
         Inches(0.5), Inches(3.75), Inches(12.33), Inches(0.55),
         size=16, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

add_text(s, "One agent. One demo. The entire vertical. Done.",
         Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.7),
         size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "The agent that built every step is the intelligence every step calls.",
         Inches(0.5), Inches(5.25), Inches(12.33), Inches(0.6),
         size=18, color=BLUE, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 7 — The Pattern Scales
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "The Pattern Scales",
         Inches(1), Inches(0.4), Inches(11.33), Inches(0.7),
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.15))

lines_left = [
    "If JiuwenClaw can build and power one experience —",
    "If it can reason about one domain —",
    "If it can power one workflow —",
]
lines_right = [
    "it can build and power many.",
    "it can reason across domains.",
    "it can power entire environments.",
]
for i, (l, r) in enumerate(zip(lines_left, lines_right)):
    y = Inches(1.5 + i * 0.85)
    add_text(s, l, Inches(0.5), y, Inches(6.5), Inches(0.75),
             size=18, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(s, r, Inches(7.0), y, Inches(6.0), Inches(0.75),
             size=18, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

# Three-layer diagram
layers = [
    ("Any Application\nTravel · Work · Creativity · Communication", WHITE,    RGBColor(0x2a, 0x2a, 0x2c)),
    ("🧠 JiuwenClaw\nArchitect · Intelligence Engine · Evolver",    WHITE,    DARK),
    ("System Resources\nBrowser · Files · APIs",                    GRAY,     RGBColor(0x1a, 0x1a, 0x1c)),
]
lw = Inches(7.5); lh = Inches(0.85)
lx = (W - lw) / 2
for i, (txt, tc, fc) in enumerate(layers):
    ly = Inches(4.0 + i * 1.05)
    add_rect(s, lx, ly, lw, lh, fc, txt, txt_size=14, txt_color=tc)
    sh = s.shapes[-1]
    sh.line.color.rgb = BLUE if i == 1 else GRAY
    sh.line.width = Pt(2 if i == 1 else 1)
    if i < len(layers) - 1:
        mid_x = lx + lw / 2
        arr = s.shapes.add_connector(1, mid_x, ly + lh, mid_x, ly + lh + Inches(0.2))
        arr.line.color.rgb = BLUE if i == 0 else GRAY
        arr.line.width = Pt(2)

add_text(s, "Intelligence not inside the app.  Beneath the app.  As both architect and engine.",
         Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.3),
         size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 8 — The openJiuwen Ecosystem: One Ecosystem. Your Choice.
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "But Here's What's Bigger.",
         Inches(1), Inches(0.4), Inches(11.33), Inches(0.7),
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.15))

add_text(s, "JiuwenClaw is part of the openJiuwen ecosystem.",
         Inches(1), Inches(1.3), Inches(11.33), Inches(0.6),
         size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "And the ecosystem gives you a choice.",
         Inches(1), Inches(1.9), Inches(11.33), Inches(0.6),
         size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# TripWise box (center top)
tw_x = Inches(5.2); tw_y = Inches(2.7); tw_w = Inches(3.0); tw_h = Inches(0.75)
add_rect(s, tw_x, tw_y, tw_w, tw_h, RGBColor(0x00, 0x40, 0x80),
         "🖥️ TripWise\n(Thin Client)", txt_size=13, txt_color=WHITE)
tw_shape = s.shapes[-1]
tw_shape.line.color.rgb = GRAY
tw_shape.line.width = Pt(1.5)

# JiuwenClaw box (left)
jc_x = Inches(1.0); jc_y = Inches(4.3); jc_w = Inches(3.8); jc_h = Inches(1.3)
add_rect(s, jc_x, jc_y, jc_w, jc_h, DARK,
         "🧠 JiuwenClaw\nLive openJiuwen Agent\nWebSocket — ready now", txt_size=14, txt_color=WHITE)
jc_eco = s.shapes[-1]
jc_eco.line.color.rgb = BLUE
jc_eco.line.width = Pt(2)

# OpenJiuwen box (right)
oj_x = Inches(8.4); oj_y = Inches(4.3); oj_w = Inches(3.8); oj_h = Inches(1.3)
add_rect(s, oj_x, oj_y, oj_w, oj_h, RGBColor(0x00, 0x3a, 0x70),
         "🔧 OpenJiuwen\nBuild Your Own Private Agent\nREST — your intelligence", txt_size=14, txt_color=WHITE)
oj_eco = s.shapes[-1]
oj_eco.line.color.rgb = TEAL
oj_eco.line.width = Pt(2)

# Arrows: TripWise → JiuwenClaw
s.shapes.add_connector(1, Inches(5.2), Inches(3.07), Inches(2.9), Inches(4.3)).line.color.rgb = BLUE
# Arrows: TripWise → OpenJiuwen
s.shapes.add_connector(1, Inches(8.2), Inches(3.07), Inches(10.3), Inches(4.3)).line.color.rgb = TEAL

# Label arrows
add_text(s, "calls for intelligence", Inches(2.3), Inches(3.6), Inches(2.8), Inches(0.4),
         size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
add_text(s, "calls for intelligence", Inches(8.3), Inches(3.6), Inches(2.8), Inches(0.4),
         size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# Labels under boxes
add_text(s, "Use the live community agent", Inches(0.8), Inches(5.65), Inches(4.3), Inches(0.45),
         size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
add_text(s, "Own your intelligence layer", Inches(8.2), Inches(5.65), Inches(4.3), Inches(0.45),
         size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

add_text(s, "Both integrations equally smooth.  Same prompts.  Same JSON back.  The thin client just calls.",
         Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.55),
         size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "The ecosystem answers. Your way.",
         Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.35),
         size=14, color=BLUE, italic=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 9 — One More Thing
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "And now…",
         Inches(1), Inches(0.8), Inches(11.33), Inches(0.8),
         size=32, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "One more thing.",
         Inches(1), Inches(1.6), Inches(11.33), Inches(0.9),
         size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "JiuwenClaw doesn't just build the product and power its intelligence.",
         Inches(1), Inches(2.8), Inches(11.33), Inches(0.6),
         size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "It evolves it.",
         Inches(1), Inches(3.4), Inches(11.33), Inches(0.8),
         size=40, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "Skills Autonomous Evolution",
         Inches(2), Inches(4.35), Inches(9.33), Inches(0.65),
         size=22, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "When JiuwenClaw encounters a failure — it captures it, studies it,\nand improves itself. Automatically. Without a single line of human code.",
         Inches(1.5), Inches(5.05), Inches(10.33), Inches(0.9),
         size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, "Architect.  Intelligence engine.  Teacher of itself.",
         Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.7),
         size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Slide 9 — The Next Chapter
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "This Is the Next Chapter",
         Inches(1), Inches(0.4), Inches(11.33), Inches(0.7),
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
divider_line(s, Inches(1.15))

bullets = [
    "The agent that builds the system is the system",
    "Experiences aren't installed — they're generated",
    "Workflows aren't configured — they're inferred",
    "Systems don't stagnate — they evolve",
    "Intelligence becomes the platform beneath everything",
]
for i, b in enumerate(bullets):
    add_text(s, f"→  {b}",
             Inches(2), Inches(1.5 + i * 0.95), Inches(9.33), Inches(0.8),
             size=21, color=WHITE if i < 4 else BLUE, bold=(i == 4),
             align=PP_ALIGN.LEFT)

# Creator / Executor / Evolver three-box at bottom
bw = Inches(3.1); bh = Inches(0.7); by = Inches(6.3)
for i, (label, col) in enumerate([
    ("🏗️  Builds Product", ORANGE),
    ("⚡  Powers Intelligence", BLUE),
    ("🔄  Evolves Itself",  GREEN),
]):
    bx2 = Inches(0.5 + i * 4.2)
    add_rect(s, bx2, by, bw, bh, col, label, txt_size=16, txt_color=WHITE, bold=True)
    if i < 2:
        mid = s.shapes.add_connector(1, bx2 + bw, by + bh/2,
                                      bx2 + bw + Inches(1.1), by + bh/2)
        mid.line.color.rgb = WHITE
        mid.line.width = Pt(2)

# ---------------------------------------------------------------------------
# Slide 10 — Closing
# ---------------------------------------------------------------------------
s = add_slide()
add_text(s, "The architect and the intelligence engine",
         Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.9),
         size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "are the same agent.",
         Inches(0.5), Inches(2.3), Inches(12.33), Inches(0.9),
         size=52, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

add_text(s, "That is the new foundation for computing.",
         Inches(0.5), Inches(3.5), Inches(12.33), Inches(0.6),
         size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "And today… we're taking the first step.",
         Inches(0.5), Inches(4.15), Inches(12.33), Inches(0.6),
         size=22, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

divider_line(s, Inches(5.2))

add_text(s, "🌐  openjiuwen.com/en",
         Inches(0.5), Inches(5.5), Inches(5.8), Inches(0.5),
         size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "💻  github.com/openJiuwen-ai",
         Inches(7.0), Inches(5.5), Inches(5.8), Inches(0.5),
         size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "JiuwenClaw  ·  openJiuwen Community  ·  2026",
         Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5),
         size=13, color=RGBColor(0x48, 0x48, 0x4a), align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = "/Users/mishka/PycharmProjects/openjiuwen/michael-agent-studio/backend/openjiuwen_studio/core/dsl_converter/demos/TripwisePython/blog/jiuwenclaw_first/BLOGPOST-KeyNote-SteveJobs.pptx"
prs.save(out)
print(f"Saved: {out}")
